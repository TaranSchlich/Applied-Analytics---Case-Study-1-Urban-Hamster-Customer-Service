from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand Colors ──────────────────────────────────────────────────────────────
PINK      = RGBColor(0xE9, 0x1E, 0x8C)
BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
BG        = RGBColor(0xF2, 0xF2, 0xF2)
GRAY_MED  = RGBColor(0xBD, 0xBD, 0xBD)
GRAY_DARK = RGBColor(0x55, 0x55, 0x55)
GRAY_LT   = RGBColor(0xE2, 0xE2, 0xE2)
PINK_LT   = RGBColor(0xF9, 0xD0, 0xEB)

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BG

def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz,
        bold=False, italic=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = "Calibri"
    r.font.size  = Pt(sz)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def mltxt(slide, lines, l, t, w, h):
    """lines = [(text, sz, bold, color, align)]"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, sz, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name  = "Calibri"
        r.font.size  = Pt(sz)
        r.font.bold  = bold
        r.font.color.rgb = color
    return tb

def slide_chrome(slide, title_text, accent_width=5.0):
    """Pink top bar + slide title + accent underline."""
    rect(slide, 0, 0, 13.33, 0.09, PINK)
    txt(slide, title_text, 0.5, 0.17, 12.3, 0.72, 25, bold=True, color=BLACK)
    rect(slide, 0.5, 0.93, accent_width, 0.045, PINK)

# ── Build Presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# COVER
# ─────────────────────────────────────────────────────────────────────────────
s0 = prs.slides.add_slide(BLANK)
set_bg(s0)
rect(s0, 0, 0,    13.33, 0.14, PINK)
rect(s0, 0, 7.36, 13.33, 0.14, PINK)

txt(s0, "URBAN HAMSTER", 0.5, 1.1, 12.33, 1.05,
    54, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
txt(s0, "C L O T H I N G", 0.5, 2.15, 12.33, 0.5,
    13, bold=False, color=GRAY_DARK, align=PP_ALIGN.CENTER)

rect(s0, 4.42, 2.8, 4.5, 0.055, PINK)

txt(s0, "Return Policy Analysis", 0.5, 2.95, 12.33, 0.85,
    36, bold=False, color=BLACK, align=PP_ALIGN.CENTER)
txt(s0, "Customer Service Team  ·  Case Study 1", 0.5, 3.85, 12.33, 0.55,
    19, color=GRAY_DARK, align=PP_ALIGN.CENTER)
txt(s0, "[ Add logo above title ]", 0.5, 0.2, 12.33, 0.7,
    11, italic=True, color=GRAY_MED, align=PP_ALIGN.CENTER)
txt(s0, "T. Schlichtmann  |  GB884  |  May 2025", 0.5, 6.7, 12.33, 0.45,
    12, color=GRAY_MED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Return Rate Overview
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
set_bg(s1)
slide_chrome(s1, "Our Return Rate Is Half the Industry Average — But Context Matters", 6.5)

# ── Horizontal bar chart (manual) ───
SCALE = 0.295   # inches per 1%

# Label header
txt(s1, "Return Rate Comparison", 0.5, 1.08, 6.5, 0.38,
    13, bold=True, color=GRAY_DARK)

# Urban Hamster bar
txt(s1, "Urban Hamster", 0.5, 1.52, 2.8, 0.32, 12, bold=True, color=BLACK)
rect(s1, 0.5, 1.85, 10 * SCALE, 0.52, PINK)
txt(s1, "10.0%", 0.5 + 10 * SCALE + 0.12, 1.85, 1.1, 0.52,
    20, bold=True, color=PINK)

# Industry bar
txt(s1, "Industry Average (Online Retail)", 0.5, 2.57, 3.5, 0.32,
    12, bold=True, color=BLACK)
rect(s1, 0.5, 2.9, 19.3 * SCALE, 0.52, GRAY_MED)
txt(s1, "19.3%", 0.5 + 19.3 * SCALE + 0.12, 2.9, 1.2, 0.52,
    20, bold=True, color=GRAY_DARK)

txt(s1, "Source: National Retail Federation, 2025", 0.5, 3.6, 6.5, 0.28,
    9, italic=True, color=GRAY_MED)

# Trend callout
rect(s1, 0.5, 4.05, 6.2, 0.05, PINK)
txt(s1, "Return Rate Trend", 0.5, 4.18, 6.2, 0.38,
    12, bold=True, color=GRAY_DARK)
txt(s1, "9.2%  →  9.5%  →  9.9%  →  9.8%  →  10.3%",
    0.5, 4.58, 6.2, 0.42, 14, bold=True, color=BLACK)
txt(s1, "2019         2020         2021         2022         2023",
    0.5, 5.0,  6.2, 0.35, 10, color=GRAY_MED)
txt(s1, "↑  Steady climb — confirms Sarah's concern about an uptick",
    0.5, 5.42, 6.2, 0.38, 11, italic=True, color=GRAY_DARK)

# ── Right insight card ───
rect(s1, 7.2, 1.05, 5.7, 5.9, GRAY_LT)
rect(s1, 7.2, 1.05, 0.07, 5.9, PINK)

txt(s1, "What This Means", 7.45, 1.15, 5.2, 0.42,
    14, bold=True, color=PINK)
mltxt(s1, [
    ("A below-average return rate sounds positive. Under a strict policy, however, "
     "it may also signal:", 12, False, GRAY_DARK, PP_ALIGN.LEFT),
    ("", 8, False, BLACK, PP_ALIGN.LEFT),
    ("  •   Hesitant customers choosing not to buy at all", 12, False, BLACK, PP_ALIGN.LEFT),
    ("  •   Frustrated returners abandoning the brand quietly", 12, False, BLACK, PP_ALIGN.LEFT),
    ("  •   Hidden lost revenue that never appears in return data", 12, False, BLACK, PP_ALIGN.LEFT),
    ("", 8, False, BLACK, PP_ALIGN.LEFT),
    ("Urban Hamster's rate is ~9 points below the industry norm. "
     "The question isn't just 'are returns low?' — it's 'why?'", 12, True, GRAY_DARK, PP_ALIGN.LEFT),
], 7.45, 1.65, 5.2, 4.5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Where Returns Are Happening
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
set_bg(s2)
slide_chrome(s2, "Returns Are Concentrated in Fit-Sensitive Categories", 5.5)

# ── Left: Category bars ───
txt(s2, "Top 3 Categories by Return Rate", 0.5, 1.08, 6.0, 0.38,
    13, bold=True, color=GRAY_DARK)

cats   = [("Pants", 11.0, PINK),
          ("Leggings", 10.7, RGBColor(0xF0, 0x7A, 0xC2)),
          ("Socks & Hosiery", 10.7, GRAY_MED)]
SCALE2 = 0.32
y = 1.55
for (name, rate, color) in cats:
    txt(s2, name, 0.5, y, 2.5, 0.3, 12, bold=False, color=BLACK)
    rect(s2, 0.5, y + 0.3, rate * SCALE2, 0.42, color)
    txt(s2, f"{rate}%", 0.5 + rate * SCALE2 + 0.12, y + 0.3, 0.9, 0.42,
        14, bold=True, color=color)
    y += 0.95

# Insight below categories
rect(s2, 0.5, 4.55, 6.0, 0.045, PINK)
mltxt(s2, [
    ("All three categories are fit-dependent — customers can't try them on. "
     "This suggests sizing guidance or product descriptions may need attention, "
     "not just the return policy itself.", 12, False, GRAY_DARK, PP_ALIGN.LEFT),
], 0.5, 4.65, 6.0, 1.5)

# ── Right: Revenue lost card ───
rect(s2, 7.1, 1.05, 5.8, 2.9, GRAY_LT)
rect(s2, 7.1, 1.05, 0.07, 2.9, PINK)
txt(s2, "Revenue Lost to Returns", 7.35, 1.15, 5.35, 0.4,
    13, bold=True, color=BLACK)
txt(s2, "$1,088,903.82", 7.35, 1.65, 5.35, 0.85,
    30, bold=True, color=PINK, align=PP_ALIGN.CENTER)
txt(s2, "in total sale price across all returned orders",
    7.35, 2.52, 5.35, 0.38, 11, italic=True,
    color=GRAY_DARK, align=PP_ALIGN.CENTER)
txt(s2, "124,877 total orders  ·  12,531 returned",
    7.35, 2.92, 5.35, 0.35, 10, color=GRAY_MED, align=PP_ALIGN.CENTER)

# Revenue context card
rect(s2, 7.1, 4.15, 5.8, 2.5, GRAY_LT)
rect(s2, 7.1, 4.15, 0.07, 2.5, PINK)
txt(s2, "This Is Just Direct Revenue Loss", 7.35, 4.25, 5.35, 0.4,
    13, bold=True, color=BLACK)
mltxt(s2, [
    ("The $1.08M does not account for:", 11, False, GRAY_DARK, PP_ALIGN.LEFT),
    ("  •   Return processing & restocking costs", 11, False, BLACK, PP_ALIGN.LEFT),
    ("  •   Lost future purchases from dissatisfied customers", 11, False, BLACK, PP_ALIGN.LEFT),
    ("  •   Brand damage from negative reviews", 11, False, BLACK, PP_ALIGN.LEFT),
    ("The true cost is higher.", 11, True, GRAY_DARK, PP_ALIGN.LEFT),
], 7.35, 4.7, 5.35, 1.8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Policy Impact & Recommendations
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
set_bg(s3)
slide_chrome(s3, "A Strict Policy Has Hidden Costs — Here's What We Recommend", 6.5)

# ── Left: Impacts ───
txt(s3, "How the Policy Is Hurting Us", 0.5, 1.08, 5.8, 0.38,
    13, bold=True, color=PINK)

impacts = [
    ("Sales",
     "Deters uncertain buyers from purchasing at all. Lost revenue that never "
     "appears in return data — the hidden cost of friction."),
    ("Reputation",
     "Difficult return experiences drive negative reviews, social media complaints, "
     "and word-of-mouth damage a refund alone can't undo."),
    ("Return Rate",
     "Strict policies don't eliminate returns. They make them harder to manage "
     "and push customers toward chargebacks and churn instead."),
]
y3 = 1.55
for (label, body) in impacts:
    rect(s3, 0.5, y3, 0.06, 0.85, PINK)
    txt(s3, label, 0.72, y3,      4.8, 0.32, 12, bold=True,  color=BLACK)
    txt(s3, body,  0.72, y3+0.3,  4.8, 0.6,  11, bold=False, color=GRAY_DARK)
    y3 += 1.05

txt(s3, "Petersen & Kumar (2009); Bower & Maxham (2012) — Journal of Marketing",
    0.5, 6.9, 6.0, 0.38, 9, italic=True, color=GRAY_MED)

# Divider
rect(s3, 6.55, 1.0, 0.045, 5.7, GRAY_MED)

# ── Right: Recommendations ───
txt(s3, "Recommended Changes", 6.8, 1.08, 6.0, 0.38,
    13, bold=True, color=PINK)

recs = [
    ("Extend return window to 30 days",
     "Current industry standard for online apparel. Competitors like Nordstrom "
     "(no time limit) and Zappos (365 days) set a high bar. 30 days is the minimum."),
    ("Offer prepaid return labels",
     "Paying for return shipping is the #1 reason shoppers abandon a brand after "
     "a bad experience. A prepaid label communicates trust and fairness."),
    ("Keep restrictions on sale & clearance only",
     "This is standard practice and widely accepted by consumers. Tighter rules "
     "on discounted items are reasonable — the main policy should not be."),
]
y4 = 1.55
for (label, body) in recs:
    rect(s3, 6.8, y4, 5.9, 0.045, PINK)
    txt(s3, label, 6.8, y4+0.07, 5.9, 0.32, 12, bold=True,  color=BLACK)
    txt(s3, body,  6.8, y4+0.38, 5.9, 0.62, 11, bold=False, color=GRAY_DARK)
    y4 += 1.05

txt(s3, "NRF, 2025; Bower & Maxham, 2012; Petersen & Kumar, 2009",
    6.8, 6.9, 6.0, 0.38, 9, italic=True, color=GRAY_MED)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/user/Applied-Analytics---Case-Study-1-Urban-Hamster-Customer-Service/GB884_Case1_ReturnPolicy_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
